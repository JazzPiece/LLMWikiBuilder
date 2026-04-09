"""
extractor.py — File text extraction.

Ported and generalised from compile_wiki.py. All extraction logic lives here;
everything is config-driven (no hardcoded extension sets — those come from WikiConfig).
"""

from __future__ import annotations

import fnmatch
import hashlib
import re
from datetime import datetime
from pathlib import Path

from .config import WikiConfig


# ---------------------------------------------------------------------------
# Maps (used for rendering, not filtering — filtering is config-driven)
# ---------------------------------------------------------------------------

FILE_TAG: dict[str, str] = {
    ".sql": "sql", ".lpd": "sql",
    ".py": "python",
    ".docx": "word", ".dotx": "word",
    ".xlsx": "excel", ".xltx": "excel",
    ".pdf": "pdf",
    ".pptx": "powerpoint",
    ".ps1": "script", ".bat": "script", ".sh": "script", ".vbs": "script",
    ".html": "html", ".htm": "html",
    ".xml": "xml",
    ".json": "json",
    ".csv": "data",
    ".ts": "typescript",
    ".js": "javascript",
    ".md": "markdown",
    ".eml": "other", ".msg": "other",
    ".vsdx": "other",
    ".vtt": "other",
}

FENCE_LANG: dict[str, str] = {
    ".sql": "sql", ".lpd": "sql",
    ".py": "python",
    ".html": "html", ".htm": "html",
    ".xml": "xml",
    ".json": "json",
    ".md": "markdown",
    ".bat": "batch",
    ".ps1": "powershell",
    ".sh": "bash",
    ".yaml": "yaml", ".yml": "yaml",
    ".js": "javascript",
    ".ts": "typescript",
    ".css": "css",
    ".vbs": "vbscript",
    ".toml": "toml",
}


# ---------------------------------------------------------------------------
# Path / name helpers
# ---------------------------------------------------------------------------

def slugify(name: str) -> str:
    """Convert a filename stem to a safe wiki article slug."""
    name = re.sub(r"[^\w\s\-]", "", name)
    name = re.sub(r"\s+", "-", name.strip())
    return name


def format_size(n: int) -> str:
    for unit in ("B", "KB", "MB", "GB"):
        if n < 1024:
            return f"{n:.0f} {unit}"
        n /= 1024
    return f"{n:.1f} GB"


def mtime_str(path: Path) -> str:
    try:
        return datetime.fromtimestamp(path.stat().st_mtime).strftime("%Y-%m-%d")
    except Exception:
        return "unknown"


def path_to_uri(p: Path) -> str:
    """Return a properly percent-encoded file:// URI. Uses Path.as_uri() which
    correctly encodes all special characters including [ ] ( ) # % spaces."""
    try:
        return Path(p).resolve().as_uri()
    except ValueError:
        # Fallback for relative paths on Windows
        return "file:///" + str(p).replace("\\", "/").replace(" ", "%20")


def file_hash(path: Path) -> str:
    h = hashlib.md5()
    try:
        with open(path, "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                h.update(chunk)
    except Exception:
        return ""
    return h.hexdigest()


def content_hash(text: str) -> str:
    """SHA-256 of content string — used as LLM cache key."""
    return hashlib.sha256(text.encode("utf-8", errors="replace")).hexdigest()


# ---------------------------------------------------------------------------
# File filtering
# ---------------------------------------------------------------------------

def should_skip_file(path: Path, cfg: WikiConfig) -> bool:
    """Return True if this file should be excluded from processing."""
    name = path.name
    for pattern in cfg.source.exclude_patterns:
        if fnmatch.fnmatch(name, pattern):
            return True
    # Size guard
    try:
        size_mb = path.stat().st_size / (1024 * 1024)
        if size_mb > cfg.source.max_file_size_mb:
            return True
    except Exception:
        pass
    return False


def should_skip_dir(dirname: str, cfg: WikiConfig) -> bool:
    """Return True if this directory name is in the exclusion list."""
    return dirname in cfg.source.exclude_folders


# ---------------------------------------------------------------------------
# Text truncation
# ---------------------------------------------------------------------------

def _truncate(text: str, max_chars: int) -> str:
    if len(text) > max_chars:
        return text[:max_chars] + f"\n\n*[… content truncated at {max_chars:,} characters]*"
    return text


# ---------------------------------------------------------------------------
# Per-format extractors
# ---------------------------------------------------------------------------

def _extract_plain(path: Path) -> tuple[str, str]:
    ext = path.suffix.lower().lstrip(".")
    label_map = {
        "sql": "SQL", "py": "Python", "html": "HTML", "htm": "HTML",
        "csv": "CSV", "xml": "XML", "json": "JSON", "md": "Markdown",
        "bat": "Batch Script", "ps1": "PowerShell", "sh": "Shell Script",
        "lpd": "Lawson LPD", "log": "Log File", "ts": "TypeScript",
        "js": "JavaScript", "css": "CSS", "yaml": "YAML", "yml": "YAML",
        "toml": "TOML", "ini": "INI", "cfg": "Config",
    }
    label = label_map.get(ext, ext.upper() if ext else "Text")
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
        return label, text
    except Exception as e:
        return label, f"*(Could not read file: {e})*"


def _extract_docx(path: Path) -> tuple[str, str]:
    try:
        from docx import Document
        doc = Document(str(path))
        lines = [para.text for para in doc.paragraphs if para.text.strip()]
        return "Word Document", "\n".join(lines)
    except ImportError:
        return "Word Document", "*(python-docx not installed — run: pip install python-docx)*"
    except Exception as e:
        return "Word Document", f"*(Extraction error: {e})*"


def _extract_xlsx(path: Path, max_rows: int = 50, max_cols: int = 20) -> tuple[str, str]:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"### Sheet: {sheet_name}\n")
            rows_seen = 0
            for row in ws.iter_rows(max_row=max_rows, max_col=max_cols, values_only=True):
                if any(cell is not None for cell in row):
                    parts.append(" | ".join(str(c) if c is not None else "" for c in row))
                    rows_seen += 1
            if rows_seen == max_rows:
                parts.append(f"*[… first {max_rows} rows shown]*")
            parts.append("")
        return "Excel Spreadsheet", "\n".join(parts)
    except ImportError:
        return "Excel Spreadsheet", "*(openpyxl not installed — run: pip install openpyxl)*"
    except Exception as e:
        return "Excel Spreadsheet", f"*(Extraction error: {e})*"


def _extract_pdf(path: Path) -> tuple[str, str]:
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(str(path)) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    texts.append(t)
        return "PDF Document", "\n\n".join(texts)
    except ImportError:
        return "PDF Document", "*(pdfplumber not installed — run: pip install pdfplumber)*"
    except Exception as e:
        return "PDF Document", f"*(Extraction error: {e})*"


def _extract_pptx(path: Path) -> tuple[str, str]:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if texts:
                slides.append(f"**Slide {i}:**\n" + "\n".join(texts))
        return "PowerPoint", "\n\n".join(slides)
    except ImportError:
        return "PowerPoint", "*(python-pptx not installed — run: pip install python-pptx)*"
    except Exception as e:
        return "PowerPoint", f"*(Extraction error: {e})*"


def _extract_eml(path: Path) -> tuple[str, str]:
    """Extract headers and body text from .eml email files."""
    import email
    from email import policy as email_policy
    try:
        raw = path.read_bytes()
        msg = email.message_from_bytes(raw, policy=email_policy.default)
        parts: list[str] = []
        for header in ("Subject", "From", "To", "Cc", "Date"):
            val = msg.get(header, "")
            if val:
                parts.append(f"**{header}:** {val}")
        parts.append("")
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    try:
                        body_bytes = part.get_payload(decode=True)
                        charset = part.get_content_charset() or "utf-8"
                        parts.append(body_bytes.decode(charset, errors="replace"))
                    except Exception:
                        pass
        else:
            if msg.get_content_type() == "text/plain":
                try:
                    body_bytes = msg.get_payload(decode=True)
                    charset = msg.get_content_charset() or "utf-8"
                    parts.append(body_bytes.decode(charset, errors="replace"))
                except Exception:
                    parts.append(str(msg.get_payload()))
        return "Email", "\n".join(parts)
    except Exception as e:
        return "Email", f"*(Extraction error: {e})*"


def _extract_vsdx(path: Path) -> tuple[str, str]:
    """Extract shape text from Visio .vsdx files (ZIP + XML)."""
    import zipfile
    import xml.etree.ElementTree as ET
    try:
        texts: list[str] = []
        with zipfile.ZipFile(str(path)) as zf:
            page_files = sorted(
                n for n in zf.namelist()
                if n.startswith("visio/pages/page") and n.endswith(".xml")
            )
            for pf in page_files:
                page_num = pf.replace("visio/pages/page", "").replace(".xml", "")
                with zf.open(pf) as f:
                    tree = ET.parse(f)
                    root = tree.getroot()
                    page_texts: list[str] = []
                    for elem in root.iter():
                        # ElementTree strips namespace prefixes — match local name
                        local = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
                        if local == "Text" and elem.text and elem.text.strip():
                            page_texts.append(elem.text.strip())
                    if page_texts:
                        texts.append(f"**Page {page_num}:**\n" + "\n".join(page_texts))
        return "Visio Diagram", "\n\n".join(texts) if texts else "*(No text content found in diagram)*"
    except Exception as e:
        return "Visio Diagram", f"*(Extraction error: {e})*"


# ---------------------------------------------------------------------------
# Public extraction entry point
# ---------------------------------------------------------------------------

def extract_text(path: Path, cfg: WikiConfig) -> tuple[str, str]:
    """
    Return (file_type_label, extracted_text).

    The text is NOT truncated here — truncation happens in article.py based
    on config, so LLM backends receive full text for summarization.
    Callers that embed raw content in wiki articles should truncate separately.
    """
    ext = path.suffix.lower()
    text_exts = cfg.text_extensions_set()
    rich_exts = cfg.rich_extensions_set()

    if ext in text_exts:
        return _extract_plain(path)
    if ext in (".docx", ".dotx"):
        return _extract_docx(path)
    if ext in (".xlsx", ".xltx"):
        return _extract_xlsx(path)
    if ext == ".pdf":
        return _extract_pdf(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext in (".eml", ".msg"):
        return _extract_eml(path)
    if ext == ".vsdx":
        return _extract_vsdx(path)
    if ext in rich_exts:
        # Fallback for any rich ext without a dedicated extractor
        size = format_size(path.stat().st_size) if path.exists() else "unknown"
        return ext.lstrip(".").upper(), f"*(File size: {size} — no text extractor for {ext})*"

    size = format_size(path.stat().st_size) if path.exists() else "unknown"
    return "Binary/Other", f"*(Binary file — no text extraction. Size: {size})*"


def get_file_tag(path: Path) -> str:
    return FILE_TAG.get(path.suffix.lower(), "other")


def get_fence_lang(path: Path) -> str:
    return FENCE_LANG.get(path.suffix.lower(), "")


# ---------------------------------------------------------------------------
# Chunking for large files
# ---------------------------------------------------------------------------

def chunk_content(content: str, max_chars: int, overlap: int) -> list[str]:
    """
    Split content into overlapping windows, preferring paragraph boundaries.
    Returns a single-element list if content fits within max_chars.
    """
    if len(content) <= max_chars:
        return [content]

    chunks: list[str] = []
    start = 0
    while start < len(content):
        end = start + max_chars
        if end >= len(content):
            chunks.append(content[start:])
            break
        # Try to split on paragraph boundary (\n\n)
        split_pos = content.rfind("\n\n", start, end)
        if split_pos == -1 or split_pos <= start:
            # Fall back to newline
            split_pos = content.rfind("\n", start, end)
        if split_pos == -1 or split_pos <= start:
            split_pos = end
        chunks.append(content[start:split_pos])
        start = max(split_pos - overlap, start + 1)

    return chunks
